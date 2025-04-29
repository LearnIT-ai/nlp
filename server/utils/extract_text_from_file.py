import win32com.client
from fastapi import  UploadFile
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import json
import os
import tempfile


class ExtractTextFromFile:
    def pdf(file: UploadFile) -> str:
        text = ""
        try:
            with fitz.open(stream=file.file.read(), filetype="pdf") as doc:
                for page in doc:
                    text += page.get_text("text") + "\n"
        except Exception as e:
            raise Exception("Problem during extracting text from pdf file") 
        return text.strip()

    def docx(file) -> str:
        try:
            file.file.seek(0)  
            doc = Document(file.file)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text.strip()
        except Exception as e:
            raise Exception(f"Error during extracting text from .docx: {str(e)}")

    def doc(file) -> str:
        try:
            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(file.file.name) 
            text = doc.Content.Text
            doc.Close()
            word.Quit()
            return text.strip()
        except Exception as e:
            raise Exception(f"Error during extracting text from .doc: {str(e)}")
        
    def pptx(file: UploadFile) -> str:
        try:
            file.file.seek(0)
            prs = Presentation(file.file)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text).strip()
        except Exception as e:
            raise Exception(f"Error during extracting text from .pptx: {str(e)}")
    
    def ppt(file: UploadFile) -> str:
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".ppt") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name

            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(tmp_path, WithWindow=False)
            text = []
            for slide in presentation.Slides:
                for shape in slide.Shapes:
                    if shape.HasTextFrame and shape.TextFrame.HasText:
                        text.append(shape.TextFrame.TextRange.Text)
            presentation.Close()
            powerpoint.Quit()
            return "\n".join(text).strip()
        except Exception as e:
            raise Exception(f"Error during extracting text from .ppt: {str(e)}")
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.remove(tmp_path)


    def py(file: UploadFile) -> str:
        try:
            file.file.seek(0)
            content = file.file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8')
            return content.strip()
        except Exception as e:
            raise Exception(f"Error during extracting text from .py: {str(e)}")

    def ipynb(file: UploadFile) -> str:
        try:
            file.file.seek(0)
            data = json.load(file.file)
            text = []
            for cell in data.get("cells", []):
                if "source" in cell:
                    text.append("".join(cell["source"]))
            return "\n".join(text).strip()
        except Exception as e:
            raise Exception(f"Error during extracting text from .ipynb: {str(e)}")

    def multi_file(file,ext) -> str:
        if ext == "pdf":
            return ExtractTextFromFile.pdf(file)
        elif ext == "docx":
            return ExtractTextFromFile.docx(file)
        elif ext == "doc":
            return ExtractTextFromFile.doc(file)
        elif ext == "pptx":
            return ExtractTextFromFile.pptx(file)
        elif ext == "ppt":
            return ExtractTextFromFile.ppt(file)
        elif ext == "py":
            return ExtractTextFromFile.py(file)
        elif ext == "ipynb":
            return ExtractTextFromFile.ipynb(file)
        else:
            raise Exception("Unsupported file extension")